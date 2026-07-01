#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""Document parsers for enterprise knowledge-base ingestion."""
import asyncio
import io
import logging
from dataclasses import dataclass
from typing import List, Optional, Tuple

logger = logging.getLogger(__name__)


@dataclass
class PageBlock:
    page_no: int
    heading_path: str
    text: str


@dataclass
class ParseResult:
    blocks: List[PageBlock]
    page_count: int = 0
    warning: Optional[str] = None


def decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "gb18030", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


async def parse_document(
        file_bytes: bytes,
        mime: str,
        filename: str = "",
        *,
        max_bytes: int,
        max_pages: int,
        ocr_mode: str = "off",
        ocr_max_pages: int = 0,
) -> ParseResult:
    data = file_bytes
    warning_parts: list[str] = []
    if len(data) > max_bytes:
        data = data[:max_bytes]
        warning_parts.append(f"文件超过 {max_bytes} 字节，已截断解析")
    result = await asyncio.to_thread(
        _parse_document_sync,
        data,
        mime,
        filename,
        max_pages,
    )
    if result.warning:
        warning_parts.append(result.warning)
    if ocr_mode != "off" and ocr_max_pages <= 0:
        warning_parts.append("OCR 未执行：ocr.max_pages 为 0")
    if warning_parts:
        result.warning = "；".join(warning_parts)
    return result


def _parse_document_sync(data: bytes, mime: str, filename: str, max_pages: int) -> ParseResult:
    lower_name = (filename or "").lower()
    mime = mime or ""
    if mime == "application/pdf" or lower_name.endswith(".pdf"):
        return _parse_pdf(data, max_pages)
    if lower_name.endswith(".docx") or mime.endswith("wordprocessingml.document"):
        return _parse_docx(data)
    if lower_name.endswith(".pptx") or mime.endswith("presentationml.presentation"):
        return _parse_pptx(data)
    if lower_name.endswith((".xlsx", ".xlsm")) or mime.endswith("spreadsheetml.sheet"):
        return _parse_xlsx(data, max_pages)
    if lower_name.endswith((".md", ".txt", ".csv", ".json")) or mime.startswith("text/"):
        text = decode_text(data)
        return _text_to_blocks(text, heading_path=filename or "文本", max_pages=max_pages)
    return _text_to_blocks(decode_text(data), heading_path=filename or "文档", max_pages=max_pages)


def _parse_pdf(data: bytes, max_pages: int) -> ParseResult:
    try:
        import fitz  # PyMuPDF
    except Exception:
        logger.warning("PyMuPDF 不可用，回退到 pypdf")
        return _parse_pdf_pypdf(data, max_pages)

    blocks: list[PageBlock] = []
    page_count = 0
    warning = None
    with fitz.open(stream=data, filetype="pdf") as doc:
        page_count = len(doc)
        limit = min(page_count, max_pages)
        if page_count > max_pages:
            warning = f"PDF 共 {page_count} 页，仅解析前 {max_pages} 页"
        for idx in range(limit):
            text = doc[idx].get_text("text").strip()
            blocks.append(PageBlock(page_no=idx + 1, heading_path=f"Page {idx + 1}", text=text))
    return ParseResult(blocks=blocks, page_count=page_count, warning=warning)


def _parse_pdf_pypdf(data: bytes, max_pages: int) -> ParseResult:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    page_count = len(reader.pages)
    warning = f"PDF 共 {page_count} 页，仅解析前 {max_pages} 页" if page_count > max_pages else None
    blocks = [
        PageBlock(
            page_no=idx + 1,
            heading_path=f"Page {idx + 1}",
            text=(page.extract_text() or "").strip(),
        )
        for idx, page in enumerate(reader.pages[:max_pages])
    ]
    return ParseResult(blocks=blocks, page_count=page_count, warning=warning)


def _parse_docx(data: bytes) -> ParseResult:
    from docx import Document

    doc = Document(io.BytesIO(data))
    blocks: list[PageBlock] = []
    current_heading = "正文"
    parts: list[str] = []
    ordinal = 1
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading"):
            if parts:
                blocks.append(PageBlock(page_no=ordinal, heading_path=current_heading, text="\n".join(parts)))
                ordinal += 1
                parts = []
            current_heading = text
        else:
            parts.append(text)
    if parts:
        blocks.append(PageBlock(page_no=ordinal, heading_path=current_heading, text="\n".join(parts)))
    return ParseResult(blocks=blocks, page_count=max(1, len(blocks)))


def _parse_pptx(data: bytes) -> ParseResult:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    blocks: list[PageBlock] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    texts.append(text)
        blocks.append(PageBlock(page_no=idx, heading_path=f"Slide {idx}", text="\n".join(texts)))
    return ParseResult(blocks=blocks, page_count=len(blocks))


def _parse_xlsx(data: bytes, max_pages: int) -> ParseResult:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[PageBlock] = []
    sheet_names = workbook.sheetnames[:max_pages]
    warning = f"Excel 共 {len(workbook.sheetnames)} 个工作表，仅解析前 {max_pages} 个" if len(workbook.sheetnames) > max_pages else None
    for idx, sheet_name in enumerate(sheet_names, start=1):
        sheet = workbook[sheet_name]
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                rows.append("\t".join(values))
        blocks.append(PageBlock(page_no=idx, heading_path=sheet_name, text="\n".join(rows)))
    return ParseResult(blocks=blocks, page_count=len(workbook.sheetnames), warning=warning)


def _text_to_blocks(text: str, heading_path: str, max_pages: int) -> ParseResult:
    lines = text.splitlines()
    blocks: list[PageBlock] = []
    current_heading = heading_path
    parts: list[str] = []
    ordinal = 1
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("#"):
            if parts:
                blocks.append(PageBlock(page_no=ordinal, heading_path=current_heading, text="\n".join(parts)))
                ordinal += 1
                parts = []
                if ordinal > max_pages:
                    break
            current_heading = stripped.lstrip("#").strip() or heading_path
            continue
        if stripped:
            parts.append(stripped)
    if parts and ordinal <= max_pages:
        blocks.append(PageBlock(page_no=ordinal, heading_path=current_heading, text="\n".join(parts)))
    if not blocks and text.strip():
        blocks.append(PageBlock(page_no=1, heading_path=heading_path, text=text.strip()))
    warning = f"文本分段超过 {max_pages} 段，已截断" if ordinal > max_pages else None
    return ParseResult(blocks=blocks, page_count=len(blocks), warning=warning)
